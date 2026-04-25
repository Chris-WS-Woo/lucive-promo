from __future__ import annotations

import csv
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
PPTX_PATH = PRESENTATION_DIR / "lucive-openbeta-plan-v1.pptx"
CSV_PATH = PRESENTATION_DIR / "lucive-openbeta-execution-table-v1.csv"
MD_PATH = PRESENTATION_DIR / "lucive-openbeta-launch-plan-v1.md"


@dataclass
class ValidationPoint:
    theme: str
    hypothesis: str
    kpis: str
    success: str
    owner: str


@dataclass
class ChecklistItem:
    phase: str
    area: str
    subarea: str
    task: str
    priority: str
    owner: str
    due: str
    completion: str
    check_point: str
    risk: str
    dependency: str = ""
    notes: str = ""


VALIDATION_POINTS = [
    ValidationPoint(
        "관계 지속성 + 서사 견인력",
        "캐릭터와의 관계, 스토리가 복귀 이유를 만든다.",
        "캐릭터 친밀도, 평균 세션 타임, 복귀 이유 설문/FGI",
        "친밀도 긍정 추이와 평균 세션 타임 유지, 설문/FGI에서 관계와 스토리가 복귀 이유로 반복 확인",
        "AI Engine Team + Alpha Test PO",
    ),
    ValidationPoint(
        "현재 수준의 UX 수용 가능성",
        "현재 수준의 UX가 유저들의 재방문을 방해하지 않을 정도다.",
        "리텐션(D1/D7 중심), UX 마찰 설문/FGI",
        "리텐션이 유지되고 UX 마찰이 재방문 차단의 주원인으로 나타나지 않음",
        "LUCIVE Product Team + Alpha Test PO",
    ),
]


CHECKLIST = [
    ChecklistItem("준비", "의사결정", "알파 목적", "알파 목적과 성공 정의 확정", "Blocker", "Alpha Test PO", "2026-04-23", "감정적 연결, 자발적 복귀, 운영 가능성, KPI 측정 범위가 1페이지로 승인됨", "4/23 킥오프", "목적이 흐리면 전체 일정이 흔들림"),
    ChecklistItem("준비", "의사결정", "기본 일정", "5월 1일 채널 오픈 / 캠페인 시작 기준 일정 확정", "Blocker", "Alpha Test PO", "2026-04-23", "4/28 퍼널 검증, 5/1 채널/캠페인 오픈, 5/8 vs 5/15 모집 마감 게이트, 5/22 권장 런칭 일정 확정", "4/23 킥오프", "런칭 게이트가 모호하면 No-Go"),
    ChecklistItem("준비", "제품", "LP 구성", "핵심 LP 3개 확정", "Blocker", "AI Engine Team", "2026-04-28", "LP 3개의 이름, 역할, 감정 포지션, 우선순위 확정", "4/28 준비 체크업", "LP 구성이 흔들리면 제품 메시지와 콘텐츠 로드맵이 모두 흔들림"),
    ChecklistItem("준비", "제품", "LP 구성", "LP별 감정 훅과 차별 포인트 정의", "Critical", "AI Engine Team", "2026-04-28", "각 LP의 판타지, 갈등, 복귀 이유 문서화", "4/28 준비 체크업", "없으면 검증력이 약해짐"),
    ChecklistItem("준비", "콘텐츠", "스토리 구조", "스토리 15개 목록과 구조 확정", "Blocker", "AI Engine Team", "2026-04-28", "15개 스토리의 LP 연결, 유형, 우선순위 확정", "4/28 준비 체크업", "복귀 콘텐츠 계획이 없으면 No-Go"),
    ChecklistItem("준비", "콘텐츠", "스토리라인", "1차 스토리라인 5개 완성", "Blocker", "AI Engine Team", "2026-04-30", "첫 5개 스토리가 커뮤니티 티징과 실제 플레이에 모두 사용 가능", "5/01 채널/캠페인 오픈 직전", "커뮤니티 예고가 실체 없는 상태가 됨"),
    ChecklistItem("준비", "콘텐츠", "스토리라인", "2차 스토리라인 5개 추가 완성", "Blocker", "AI Engine Team", "2026-05-07", "누적 10개 이상의 스토리가 QA 통과 후 플레이 가능", "W1 리뷰", "복귀 루프가 얕아짐"),
    ChecklistItem("준비", "콘텐츠", "스토리라인", "3차 스토리라인 5개 추가 완성", "Critical", "AI Engine Team", "2026-05-14", "누적 15개 스토리의 알파 버전이 준비됨", "런칭 직전 Go/No-Go", "런칭 볼륨이 약해지면 Conditional Go"),
    ChecklistItem("준비", "콘텐츠", "몰입 자산", "LP 3개 세계관 요약본 준비", "Blocker", "AI Engine Team", "2026-04-28", "LP 1, 2, 3 각각 1페이지 세계관 요약 존재", "4/28 준비 체크업", "커뮤니티 스토리라인 재료 부족"),
    ChecklistItem("준비", "콘텐츠", "몰입 자산", "LP 3개 채팅용 이미지 최소 세트 준비", "Blocker", "AI Engine Team", "2026-05-01", "각 LP별 티징 가능한 이미지 팩 준비", "5/01 채널/캠페인 오픈", "몰입감 약화"),
    ChecklistItem("준비", "콘텐츠", "몰입 자산", "LP 3개 보이스 샘플 최소 세트 준비", "Critical", "AI Engine Team", "2026-05-07", "각 LP별 인트로 및 핵심 감정 비트 샘플 확보", "W1 리뷰", "감정 연결 약화"),
    ChecklistItem("준비", "콘텐츠", "커뮤니티 재료", "일자별 커뮤니티 공개용 스포일러 안전 문구 준비", "Critical", "AI Engine Team", "2026-04-29", "매일 게시 가능한 대사, 설정 조각, 스토리 훅 승인 완료", "5/01 채널/캠페인 오픈", "콘텐츠 캘린더 실행력 저하"),
    ChecklistItem("준비", "마케팅", "메시지", "알파 핵심 캠페인 메시지 확정", "Blocker", "Growth Marketing Team", "2026-04-25", "KR/JP/US 타깃용 기본 메시지 세트 확정", "4/25 준비 체크업", "획득 메시지가 흔들리면 No-Go"),
    ChecklistItem("준비", "마케팅", "자산", "바이럴 가능 자산 3종 제작", "Critical", "Growth Marketing Team", "2026-04-30", "캐릭터 티저, 스토리 티저, 홍보 스토리라인 자산 준비", "5/01 채널/캠페인 오픈", "모집 효율 저하"),
    ChecklistItem("준비", "마케팅", "랜딩 퍼널", "랜딩 및 사전등록 E2E 검증", "Blocker", "LUCIVE Product Team", "2026-04-28", "모바일/데스크톱에서 랜딩, 등록, 리다이렉트, 기본 로깅 확인", "4/28 준비 체크업", "모집 퍼널이 깨지면 No-Go"),
    ChecklistItem("준비", "마케팅", "스크리닝 퍼널", "설문 응답 저장과 스크리닝 동선 검증", "Blocker", "Alpha Test PO", "2026-04-28", "응답이 선발용 필드와 함께 안정적으로 저장됨", "4/28 준비 체크업", "스크리닝 데이터 불안정"),
    ChecklistItem("준비", "마케팅", "유료 집행", "X / Meta / Google SA 세팅 완료", "Blocker", "Growth Marketing Team", "2026-04-30", "채널, 크리에이티브, 카피, UTM, 예산표 준비", "5/01 채널/캠페인 오픈", "유료 집행 지연"),
    ChecklistItem("준비", "마케팅", "유료 집행", "유료 마케팅 시작", "Blocker", "Growth Marketing Team", "2026-05-01", "캠페인 활성화 및 기본 트래킹 검증", "5/01 채널/캠페인 오픈", "모집 시작 지연"),
    ChecklistItem("준비", "커뮤니티", "공식 채널", "KR 카카오 오픈채팅 오픈", "Blocker", "Growth Marketing Team", "2026-04-29", "소개글, 규칙, 피드백 경로, 고정 공지 포함해 오픈", "5/01 채널/캠페인 오픈", "KR 유입이 머물 곳 없음"),
    ChecklistItem("준비", "커뮤니티", "공식 채널", "JP LINE 채널 오픈", "Blocker", "Growth Marketing Team", "2026-04-30", "입장 링크, 웰컴 메시지, 규칙 검증", "5/01 채널/캠페인 오픈", "JP 유입이 머물 곳 없음"),
    ChecklistItem("준비", "커뮤니티", "공식 채널", "US Discord 서버 오픈", "Blocker", "Growth Marketing Team", "2026-04-30", "채널 구조, 소개, 권한, FAQ, 피드백 경로 준비", "5/01 채널/캠페인 오픈", "US 유입이 머물 곳 없음"),
    ChecklistItem("준비", "커뮤니티", "운영 가이드", "커뮤니티 운영 가이드 공개", "Blocker", "Growth Marketing Team", "2026-04-28", "톤, 응답 규칙, escalation, 스포일러 정책, 업데이트 주기 문서화", "4/28 준비 체크업", "운영 규칙 없이 채널 오픈"),
    ChecklistItem("준비", "커뮤니티", "콘텐츠 로드맵", "5/1~5/7 커뮤니티 콘텐츠 로드맵 발행", "Blocker", "Growth Marketing Team", "2026-04-29", "LP 공개, 스토리 티저, 개발 소식, 모집 리마인드 일자별 캘린더 완성", "5/01 채널/캠페인 오픈", "채널이 비어 있음"),
    ChecklistItem("준비", "운영 인프라", "페이지 운영", "배포/롤백/수정 오너 프로세스 정의", "Blocker", "LUCIVE Product Team", "2026-04-24", "배포 담당, 롤백 기준, 긴급 수정 경로 문서화", "4/25 준비 체크업", "랜딩을 안전하게 운영 못 함"),
    ChecklistItem("준비", "운영 인프라", "링크 QA", "전체 링크 QA 시트 준비", "Blocker", "LUCIVE Product Team", "2026-04-28", "CTA, 이메일, 설문, 채널 링크를 한 시트에서 검증 가능", "4/28 준비 체크업", "유저 경로 링크 불안정"),
    ChecklistItem("준비", "운영 인프라", "Integration", "공식사이트 Integration 최종 체크", "Blocker", "LUCIVE Product Team", "2026-04-30", "공식사이트에서 랜딩, 사전등록, 설문, 채널 링크가 실제 운영 환경 기준으로 최종 연결 확인", "5/01 채널/캠페인 오픈 직전", "오픈 직후 전환 누락"),
    ChecklistItem("준비", "DB / 운영", "원장", "알파 마스터 로스터 생성", "Blocker", "Alpha Test PO", "2026-04-29", "모집 상태, 시장, 유입, 대기, 커뮤니티 입장이 단일 원장에서 추적됨", "5/01 채널/캠페인 오픈", "대상 상태 추적 불가"),
    ChecklistItem("준비", "DB / 운영", "필드 설계", "로스터 필드와 상태값 체계 정의", "Blocker", "Alpha Test PO", "2026-04-29", "이메일, 시장, 유입, 설문, 선발, 커뮤니티, FGI, CS 필드 잠금", "5/01 채널/캠페인 오픈", "운영 상태 모호"),
    ChecklistItem("준비", "본 프로덕트 운영", "버그 접수", "버그 리포트 채널, 태그, 심각도 기준 표준화", "Blocker", "LUCIVE Product Team", "2026-05-10", "단일 버그 접수 경로와 상태/심각도 태그 체계 완성", "런칭 직전 운영 체크", "피드백이 채팅 속에 묻힘"),
    ChecklistItem("준비", "본 프로덕트 운영", "유저 응대", "버그/이슈/핫픽스 회신 템플릿 준비", "Critical", "LUCIVE Product Team", "2026-05-11", "접수, 확인중, 수정완료, 알려진 이슈 템플릿 준비", "런칭 직전 운영 체크", "신뢰 리스크 증가"),
    ChecklistItem("준비", "리서치 / FGI", "후속 설문", "FGI / UT 후속 설문 준비", "Critical", "Alpha Test PO", "2026-05-08", "복귀 이유, 몰입, 마찰, LP 선호 질문 포함", "W1 리뷰 전", "정성 학습 약화"),
    ChecklistItem("준비", "리서치 / FGI", "대상 관리", "FGI 후보 상태값과 선별 규칙 추가", "Critical", "Alpha Test PO", "2026-05-08", "잔존, 이탈, 추천 유입, FGI 후보가 로스터에서 구분됨", "W1 리뷰 전", "인터뷰 모집 어려움"),
    ChecklistItem("준비", "측정", "KPI 정의", "KPI 이벤트와 산식 정의", "Blocker", "Alpha Test PO", "2026-04-29", "D7, D30, 같은 LP 재방문, 온보딩 완료, 세션 길이, 추천 관련 KPI 정의", "5/01 채널/캠페인 오픈", "무엇도 증명 못 하면 No-Go"),
    ChecklistItem("준비", "측정", "계측", "이벤트 로깅 구현 및 QA", "Blocker", "LUCIVE Product Team", "2026-05-05", "테스트 환경에서 이벤트가 정확히 발화", "W1 리뷰", "KPI 수집 불안정"),
    ChecklistItem("준비", "측정", "리포팅", "KPI 트래커와 일일 리포트 시트 구축", "Critical", "Alpha Test PO", "2026-05-08", "매일 볼 수 있는 KPI 시트와 공유 루틴 준비", "W1 리뷰", "일일 통제력 약화"),
    ChecklistItem("준비", "마케팅", "모집 게이트", "모집 마감일 5/8 vs 5/15 결정", "Blocker", "Alpha Test PO + Growth Marketing Team", "2026-05-06", "퍼널 품질, 모집 속도, 스크리닝 통과율 기준으로 마감일 최종 확정", "5/06 운영 체크업", "모집 전략과 선발 일정 혼선"),
    ChecklistItem("준비", "런칭 운영", "첫 72시간", "첫 72시간 운영 커버리지와 모니터링 역할 지정", "Blocker", "LUCIVE Product Team", "2026-05-12", "런칭일, 버그 triage, 공지, escalation 담당자 명시", "런칭 직전 Go/No-Go", "첫 주말 오너 없음"),
    ChecklistItem("준비", "런칭 운영", "리허설", "전체 런칭 리허설 실행", "Blocker", "LUCIVE Product Team", "2026-05-14", "초대, 입장, 플레이, 피드백, 이슈 공지, KPI 체크까지 드라이런 완료", "5/14 Go/No-Go", "드라이런에서 큰 결함 발견"),
    ChecklistItem("준비", "런칭 운영", "게이트 판단", "Go / No-Go 신호등 판단 공표", "Blocker", "Alpha Test PO", "2026-05-14", "Blocker/Critical 상태 기반 최종 판단 공표", "5/14 Go/No-Go", "공식 런칭 게이트"),
    ChecklistItem("운영", "알파 주간 운영", "Week 1", "운영 정책 확정", "Blocker", "Alpha Test PO", "2026-05-20", "일일 체크인, 커뮤니티 공지, 버그 응대 SLA, 핫픽스 기준, 피드백 라우팅 정책 확정", "런칭 직후 D+1", "첫 주 운영 원칙 부재"),
    ChecklistItem("운영", "알파 주간 운영", "Week 1", "일일 체크업 운영", "Blocker", "Alpha Test PO", "2026-05-23", "D+1부터 D+7까지 KPI, 커뮤니티, 버그 현황 일일 점검", "W1 일일 체크", "이탈과 장애를 초기에 못 잡음"),
    ChecklistItem("운영", "알파 주간 운영", "Week 1", "중간 점검 회의", "Blocker", "Alpha Test PO", "2026-05-27", "런칭 후 첫 5일 데이터와 이슈를 기준으로 범위/우선순위 재정렬", "W1 중간 리뷰", "문제 우선순위 미정"),
    ChecklistItem("운영", "리서치 / FGI", "Week 1", "1차 FGI 대상 확정", "Critical", "Alpha Test PO", "2026-05-27", "잔존 유저, 이탈 유저, 추천 유입 유저를 포함한 1차 FGI 대상군 확정", "W1 중간 리뷰", "정성 표본 편향"),
    ChecklistItem("운영", "리서치 / FGI", "Week 1", "1차 FGI 진행", "Critical", "Alpha Test PO", "2026-05-29", "최소 3명 이상 인터뷰/FGI 진행 및 기록 완료", "W1 종료 리뷰", "초기 정성 증거 부족"),
    ChecklistItem("운영", "커뮤니티", "Week 2", "콘텐츠 정책 확정", "Critical", "Growth Marketing Team", "2026-05-28", "Week 2용 개발 로그, 반영 포인트, 알려진 이슈, 티저 콘텐츠 계획 확정", "W2 시작", "커뮤니티 동력 저하"),
    ChecklistItem("운영", "본 프로덕트 운영", "Week 2", "알려진 이슈 보드 운영 시작", "Critical", "LUCIVE Product Team", "2026-05-28", "알려진 이슈, 수정중, 수정완료 상태를 외부 공지 가능한 포맷으로 운영", "W2 시작", "버그 신뢰도 하락"),
    ChecklistItem("운영", "알파 주간 운영", "Week 2", "중간 KPI 체크업", "Blocker", "Alpha Test PO", "2026-06-03", "D7, 온보딩, 세션 길이, 같은 LP 재방문 초기 추세 리뷰", "W2 중간 리뷰", "핵심 검증 상태를 제때 못 판단함"),
    ChecklistItem("운영", "리서치 / FGI", "Week 2", "2차 FGI 대상 확정", "Critical", "Alpha Test PO", "2026-06-03", "잔존/이탈/커뮤니티 활동도 기준으로 2차 FGI 대상군 업데이트", "W2 중간 리뷰", "FGI 표본 관리 실패"),
    ChecklistItem("운영", "리서치 / FGI", "Week 2", "2차 FGI 진행", "Critical", "Alpha Test PO", "2026-06-05", "최소 3명 이상 2차 FGI 완료 및 주요 인용문 정리", "W2 종료 리뷰", "후반 증거 부족"),
    ChecklistItem("운영", "커뮤니티", "Week 3", "커뮤니티 운영 정책 업데이트", "Critical", "Growth Marketing Team", "2026-06-04", "반영된 개선점, 향후 계획, 커뮤니티 기대치 관리 정책 재정리", "W3 시작", "기대치 불일치"),
    ChecklistItem("운영", "본 프로덕트 운영", "Week 3", "버그 / 피드백 triage 재정렬", "Critical", "LUCIVE Product Team", "2026-06-04", "누적 버그와 피드백을 제품, 콘텐츠, 운영으로 분류 재정렬", "W3 시작", "핫픽스 우선순위 혼선"),
    ChecklistItem("운영", "알파 주간 운영", "Week 3", "중간 점검", "Blocker", "Alpha Test PO", "2026-06-10", "리텐션과 커뮤니티 분위기, FGI 인사이트를 종합한 중간 리포트 작성", "W3 중간 리뷰", "중간 의사결정 지연"),
    ChecklistItem("운영", "커뮤니티", "Week 4", "커뮤니티 마감 커뮤니케이션 준비", "Critical", "Growth Marketing Team", "2026-06-11", "알파 종료 리마인드, 감사 메시지, 후속 업데이트 예고 준비", "W4 시작", "종료 커뮤니케이션 혼선"),
    ChecklistItem("운영", "리서치 / FGI", "Week 4", "최종 FGI / UT 보강 인터뷰", "Medium", "Alpha Test PO", "2026-06-12", "보강이 필요한 세그먼트 대상 추가 인터뷰 완료", "W4 시작", "일부 세그먼트 인사이트 부족"),
    ChecklistItem("운영", "알파 주간 운영", "Week 4", "최종 체크업", "Blocker", "Alpha Test PO", "2026-06-17", "4주 누적 KPI, 커뮤니티 운영, 버그 추세, FGI 인사이트 종합", "W4 최종 리뷰", "알파 결론 정리 실패"),
    ChecklistItem("운영", "알파 종료", "마감 커뮤니케이션", "알파 종료 공지 및 후속 공지 발송", "Blocker", "Growth Marketing Team", "2026-06-19", "커뮤니티, 이메일, 내부 로스터 기준 종료/감사/후속 안내 완료", "알파 종료일", "종료 메시지 혼선"),
    ChecklistItem("운영", "알파 종료", "최종 보고", "최종 알파 결과 리포트 작성", "Blocker", "Alpha Test PO", "2026-06-22", "검증 포인트별 KPI 결과, FGI 요약, 운영 회고, 다음 액션 정리", "종료 후 D+3", "알파가 기록으로 남지 않음"),
]


DEPENDENCY_MAP = {
    "LP별 감정 훅과 차별 포인트 정의": "핵심 LP 3개 확정",
    "스토리 15개 목록과 구조 확정": "핵심 LP 3개 확정",
    "1차 스토리라인 5개 완성": "핵심 LP 3개 확정, 스토리 15개 목록과 구조 확정",
    "2차 스토리라인 5개 추가 완성": "1차 스토리라인 5개 완성",
    "3차 스토리라인 5개 추가 완성": "2차 스토리라인 5개 추가 완성",
    "LP 3개 세계관 요약본 준비": "핵심 LP 3개 확정",
    "LP 3개 채팅용 이미지 최소 세트 준비": "핵심 LP 3개 확정, LP 3개 세계관 요약본 준비",
    "LP 3개 보이스 샘플 최소 세트 준비": "핵심 LP 3개 확정, LP별 감정 훅과 차별 포인트 정의",
    "일자별 커뮤니티 공개용 스포일러 안전 문구 준비": "핵심 LP 3개 확정, 스토리 15개 목록과 구조 확정, LP 3개 세계관 요약본 준비",
    "알파 핵심 캠페인 메시지 확정": "알파 목적과 성공 정의 확정, 핵심 LP 3개 확정",
    "바이럴 가능 자산 3종 제작": "알파 핵심 캠페인 메시지 확정, 핵심 LP 3개 확정, 1차 스토리라인 5개 완성, LP 3개 채팅용 이미지 최소 세트 준비",
    "랜딩 및 사전등록 E2E 검증": "배포/롤백/수정 오너 프로세스 정의, 전체 링크 QA 시트 준비",
    "설문 응답 저장과 스크리닝 동선 검증": "알파 마스터 로스터 생성, 로스터 필드와 상태값 체계 정의",
    "X / Meta / Google SA 세팅 완료": "알파 핵심 캠페인 메시지 확정, 바이럴 가능 자산 3종 제작, 랜딩 및 사전등록 E2E 검증",
    "유료 마케팅 시작": "X / Meta / Google SA 세팅 완료, KR 카카오 오픈채팅 오픈, JP LINE 채널 오픈, US Discord 서버 오픈",
    "KR 카카오 오픈채팅 오픈": "커뮤니티 운영 가이드 공개, 5/1~5/7 커뮤니티 콘텐츠 로드맵 발행",
    "JP LINE 채널 오픈": "커뮤니티 운영 가이드 공개, 5/1~5/7 커뮤니티 콘텐츠 로드맵 발행",
    "US Discord 서버 오픈": "커뮤니티 운영 가이드 공개, 5/1~5/7 커뮤니티 콘텐츠 로드맵 발행",
    "커뮤니티 운영 가이드 공개": "알파 목적과 성공 정의 확정",
    "5/1~5/7 커뮤니티 콘텐츠 로드맵 발행": "핵심 LP 3개 확정, 일자별 커뮤니티 공개용 스포일러 안전 문구 준비, 알파 핵심 캠페인 메시지 확정",
    "전체 링크 QA 시트 준비": "랜딩 및 사전등록 E2E 검증, 설문 응답 저장과 스크리닝 동선 검증, KR 카카오 오픈채팅 오픈, JP LINE 채널 오픈, US Discord 서버 오픈",
    "공식사이트 Integration 최종 체크": "랜딩 및 사전등록 E2E 검증, 설문 응답 저장과 스크리닝 동선 검증, KR 카카오 오픈채팅 오픈, JP LINE 채널 오픈, US Discord 서버 오픈",
    "알파 마스터 로스터 생성": "알파 목적과 성공 정의 확정, 5월 1일 채널 오픈 / 캠페인 시작 기준 일정 확정",
    "로스터 필드와 상태값 체계 정의": "알파 마스터 로스터 생성",
    "버그 리포트 채널, 태그, 심각도 기준 표준화": "커뮤니티 운영 가이드 공개, 알파 마스터 로스터 생성",
    "버그/이슈/핫픽스 회신 템플릿 준비": "버그 리포트 채널, 태그, 심각도 기준 표준화",
    "FGI / UT 후속 설문 준비": "스토리 15개 목록과 구조 확정, KPI 이벤트와 산식 정의",
    "FGI 후보 상태값과 선별 규칙 추가": "알파 마스터 로스터 생성, 로스터 필드와 상태값 체계 정의",
    "KPI 이벤트와 산식 정의": "알파 목적과 성공 정의 확정, 핵심 LP 3개 확정",
    "이벤트 로깅 구현 및 QA": "KPI 이벤트와 산식 정의, 랜딩 및 사전등록 E2E 검증",
    "KPI 트래커와 일일 리포트 시트 구축": "KPI 이벤트와 산식 정의, 이벤트 로깅 구현 및 QA",
    "모집 마감일 5/8 vs 5/15 결정": "유료 마케팅 시작, KPI 트래커와 일일 리포트 시트 구축",
    "첫 72시간 운영 커버리지와 모니터링 역할 지정": "버그 리포트 채널, 태그, 심각도 기준 표준화, 버그/이슈/핫픽스 회신 템플릿 준비, KPI 트래커와 일일 리포트 시트 구축",
    "전체 런칭 리허설 실행": "랜딩 및 사전등록 E2E 검증, 설문 응답 저장과 스크리닝 동선 검증, 첫 72시간 운영 커버리지와 모니터링 역할 지정, 1차 스토리라인 5개 완성",
    "Go / No-Go 신호등 판단 공표": "전체 런칭 리허설 실행, 2차 스토리라인 5개 추가 완성, KPI 트래커와 일일 리포트 시트 구축",
    "운영 정책 확정": "Go / No-Go 신호등 판단 공표, 첫 72시간 운영 커버리지와 모니터링 역할 지정",
    "일일 체크업 운영": "운영 정책 확정, KPI 트래커와 일일 리포트 시트 구축",
    "중간 점검 회의": "일일 체크업 운영",
    "1차 FGI 대상 확정": "FGI 후보 상태값과 선별 규칙 추가, 일일 체크업 운영",
    "1차 FGI 진행": "1차 FGI 대상 확정, FGI / UT 후속 설문 준비",
    "콘텐츠 정책 확정": "5/1~5/7 커뮤니티 콘텐츠 로드맵 발행, 1차 FGI 진행",
    "알려진 이슈 보드 운영 시작": "버그 리포트 채널, 태그, 심각도 기준 표준화, 버그/이슈/핫픽스 회신 템플릿 준비",
    "중간 KPI 체크업": "일일 체크업 운영, 이벤트 로깅 구현 및 QA",
    "2차 FGI 대상 확정": "1차 FGI 진행, 중간 KPI 체크업",
    "2차 FGI 진행": "2차 FGI 대상 확정",
    "커뮤니티 운영 정책 업데이트": "콘텐츠 정책 확정, 1차 FGI 진행",
    "버그 / 피드백 triage 재정렬": "알려진 이슈 보드 운영 시작, 중간 KPI 체크업",
    "중간 점검": "중간 KPI 체크업, 2차 FGI 진행, 버그 / 피드백 triage 재정렬",
    "커뮤니티 마감 커뮤니케이션 준비": "커뮤니티 운영 정책 업데이트, 중간 점검",
    "최종 FGI / UT 보강 인터뷰": "중간 점검",
    "최종 체크업": "중간 점검, 최종 FGI / UT 보강 인터뷰",
    "알파 종료 공지 및 후속 공지 발송": "커뮤니티 마감 커뮤니케이션 준비, 최종 체크업",
    "최종 알파 결과 리포트 작성": "최종 체크업, 알파 종료 공지 및 후속 공지 발송",
}

for item in CHECKLIST:
    item.dependency = DEPENDENCY_MAP.get(item.task, item.dependency)


TITLE = RGBColor(24, 35, 64)
MUTED = RGBColor(96, 108, 132)
BG = RGBColor(246, 247, 251)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(214, 84, 54)
BLUE = RGBColor(55, 96, 178)
GREEN = RGBColor(50, 128, 83)
RED = RGBColor(191, 60, 49)


def add_title(slide, title: str, subtitle: str = "") -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TITLE
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.62), Inches(0.92), Inches(12.0), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED


def add_card(slide, left: float, top: float, width: float, height: float, heading: str, body: list[str], fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = heading
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for line in body:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11.5)
        p.font.color.rgb = WHITE
        p.space_before = Pt(4)


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = TITLE
        p.space_after = Pt(8)


def add_table(slide, left: float, top: float, width: float, height: float, headers: list[str], rows: list[list[str]], widths: list[float]) -> None:
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for idx, width_inches in enumerate(widths[: len(headers)]):
        table.columns[idx].width = Inches(width_inches)
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9.5)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG if row_idx % 2 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8.3)
            p.font.color.rgb = TITLE


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Lucive 오픈베타 검증 및 운영 계획", "2026-04-24 | 준비 단계 + 4주 운영 단계 + KPI + Go/No-Go")
    add_card(slide, 0.7, 1.45, 3.85, 1.8, "이번 알파의 목적", [
        "스토리 기반 LP가 감정적 연결과 자발적 복귀를 만드는지 검증",
        "모집-선발-온보딩-커뮤니티-피드백 루프를 운영 가능한 형태로 실제로 돌려보기",
    ], BLUE)
    add_card(slide, 4.75, 1.45, 3.85, 1.8, "반드시 측정할 것", [
        "D7 / D30 리텐션",
        "같은 LP 재방문율",
        "온보딩 완료율과 세션 길이",
        "추천 유입 품질과 FGI 증거",
    ], ACCENT)
    add_card(slide, 8.8, 1.45, 3.85, 1.8, "운영 원칙", [
        "5월 1일 전에 공식 채널 오픈",
        "알파는 4주간 운영하며 중간 체크업과 FGI를 병행",
        "Alpha Test PO가 전체 readiness와 운영 체크 게이트 오너",
    ], GREEN)
    add_bullets(slide, [
        "이번 버전은 알파가 아니라 설문을 동반한 오픈베타 운영 기준으로 정리한다.",
        "검증 포인트는 2개다: 관계 지속성 + 서사 견인력, 현재 수준의 UX 수용 가능성.",
        "준비 일정뿐 아니라 4주 운영, FGI, 커뮤니티 정책, 버그/피드백 운영 시점까지 포함한다.",
    ], 0.8, 3.8, 12.0, 2.3, 16)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "이번 오픈베타에서 검증할 것", "각 검증 포인트는 KPI와 오너, 운영 체크 시점으로 연결되어야 한다.")
    rows = [[v.theme, v.hypothesis, v.kpis, v.success, v.owner] for v in VALIDATION_POINTS]
    add_table(slide, 0.45, 1.45, 12.4, 5.6, ["검증 축", "가설", "KPI", "성공 기준", "오너"], rows, [1.4, 3.3, 2.6, 3.0, 2.1])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "준비 타임라인", "주신 일정표 기준으로 5월 1일 이전 준비와 5월 중순 런칭 게이트를 정리")
    add_bullets(slide, [
        "4/23-4/25: 오픈베타 목적, 일정, 캠페인 메시지 확정",
        "4/26-4/28: LP 3개 확정, 세계관 요약, 랜딩 QA, 사전등록 QA, 설문 QA, 온보딩 측정 구조 정리",
        "4/29-4/30: KR/JP/US 공식 채널 오픈, 운영 가이드 확정, 1차 스토리라인 5개 완성, 5/1~5/7 콘텐츠 로드맵 준비, 공식사이트 Integration 최종 체크",
        "5/01: 공식 채널 + 유료 캠페인 동시 오픈",
        "5/01-5/07: 모집 진행, 일자별 개발/커뮤니티 스토리라인 운영, 2차 스토리라인 5개 완성, 퍼널 품질 모니터링",
        "5/06: 모집 마감일을 5/8로 갈지 5/15로 늘릴지 결정",
        "5/08-5/12: 선발, 온보딩 가이드, KPI 트래커, 버그 응대 체계, FGI 후속 설문 마무리",
        "5/14: 3차 스토리라인 5개 완성, 전체 리허설과 Go/No-Go 판단",
        "5/15 or 5/22: 권장 오픈베타 런칭일",
    ], 0.8, 1.55, 12.0, 5.3, 17)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "오픈베타 4주 운영 타임라인", "운영 기간에는 주차별 체크업, 커뮤니티 정책, 버그/피드백 운영, FGI 일정이 함께 돌아가야 한다.")
    add_bullets(slide, [
        "Week 1: 런칭 직후 일일 체크업, 버그 접수/응대 안정화, 1차 FGI 대상 확정",
        "Week 2: D7 초기 지표 리뷰, 알려진 이슈 보드 운영, 1차/2차 FGI 진행",
        "Week 3: 리텐션/커뮤니티 분위기 중간 리뷰, 버그/피드백 triage 재정렬, 운영 정책 업데이트",
        "Week 4: 종료 커뮤니케이션 준비, 보강 인터뷰, 최종 KPI/FGI 종합, 종료 공지",
        "D+1, W1 중간, W2 중간, W3 중간, W4 종료 직전이 핵심 체크 포인트다.",
    ], 0.8, 1.55, 12.0, 5.3, 17)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "운영 중 꼭 챙길 것", "커뮤니티, 버그, 피드백, FGI는 일정과 함께 정책도 미리 준비돼 있어야 한다.")
    add_card(slide, 0.7, 1.5, 3.0, 4.9, "커뮤니티 운영", [
        "첫 4주 콘텐츠 정책",
        "주간 공지 / 개발 로그 / 기대치 관리",
        "알파 종료 커뮤니케이션",
    ], BLUE)
    add_card(slide, 3.95, 1.5, 3.0, 4.9, "버그 / 피드백", [
        "접수 채널 단일화",
        "심각도, 상태, 공지 방식 표준화",
        "피드백 triage 재정렬 주간 운영",
    ], ACCENT)
    add_card(slide, 7.2, 1.5, 2.7, 4.9, "FGI / 리서치", [
        "1차/2차 FGI 일정",
        "대상군 관리",
        "인용문 정리",
    ], GREEN)
    add_card(slide, 10.15, 1.5, 2.5, 4.9, "체크업", [
        "D+1",
        "W1 중간",
        "W2 중간",
        "W3 중간",
        "W4 최종",
    ], RED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "상위 체크리스트 스냅샷", "준비 단계와 운영 단계를 하나의 관리표로 합쳤다.")
    snapshot = [[item.phase, item.area, item.task, item.priority, item.owner, item.due, item.check_point] for item in CHECKLIST if item.priority in {"Blocker", "Critical"}][:16]
    add_table(slide, 0.3, 1.4, 12.8, 5.9, ["단계", "영역", "항목", "중요도", "담당팀", "기한", "체크 시점"], snapshot, [0.9, 1.2, 4.4, 0.9, 1.9, 1.2, 1.8])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Go / No-Go 디시전 포인트", "제품, 퍼널, 커뮤니티, 측정, 운영 중 하나라도 Blocker가 비면 런칭하면 안 된다.")
    go_rows = [
        ["제품 최소 사용 가능", "LP 3개 플레이 가능, 진행 불가 버그 없음", "LP 하나 이상 unusable 또는 치명 버그 지속", "제품 자체가 안 돌면 알파의 의미가 없음"],
        ["복귀를 만들 콘텐츠 확보", "플레이 가능한 스토리 10개 이상, 15개 구조 확정", "복귀 이유가 너무 얕아 가설 검증 불가", "이번 알파는 서사 견인력 검증이 핵심"],
        ["캠페인 전 채널 오픈", "KR/JP/US 공식 채널 오픈 + 첫 주 게시물 준비", "유료 유입이 빈 채널로 떨어짐", "기대감은 캠페인 전에 쌓아야 함"],
        ["모집 퍼널 정상 동작", "랜딩 -> 등록 -> 설문 -> 로스터 저장이 한 바퀴 돈다", "중간 어디선가 끊기거나 데이터가 유실됨", "모집은 되는데 선발 운영이 안 되는 상황 방지"],
        ["핵심 KPI 측정 가능", "D7, 같은 LP 재방문, 온보딩, 세션 길이 측정 가능", "런칭 후에도 핵심 지표를 못 봄", "검증형 알파 목적 상실"],
        ["첫 72시간 + 4주 운영 가능", "일일 체크업, 버그/피드백 정책, FGI 일정 존재", "첫 주말 오너 없음, 운영 정책 없음", "초기 신뢰 붕괴와 중간 판단 실패 방지"],
    ]
    add_table(slide, 0.45, 1.45, 12.4, 5.8, ["디시전 포인트", "Go 기준", "No-Go 기준", "사유"], go_rows, [2.4, 3.2, 3.0, 3.4])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "이번 주 바로 볼 것", "준비 단계 마감과 운영 단계 설계를 동시에 당겨서 봐야 한다.")
    add_card(slide, 0.7, 1.55, 3.85, 4.8, "4/28까지", [
        "LP 3개 확정",
        "랜딩/사전등록/설문 QA 완료",
        "세계관 요약과 스토리 15개 구조 확정",
    ], TITLE)
    add_card(slide, 4.75, 1.55, 3.85, 4.8, "5/1까지", [
        "1차 스토리라인 5개 완성",
        "KR/JP/US 공식 채널 오픈",
        "5/1~5/7 커뮤니티 콘텐츠 로드맵 게시 가능 상태",
    ], ACCENT)
    add_card(slide, 8.8, 1.55, 3.85, 4.8, "런칭 전까지", [
        "2차/3차 스토리라인 완성",
        "버그/피드백 운영 정책 확정",
        "4주 운영 체크업 / FGI 일정표 준비",
    ], GREEN)

    prs.save(PPTX_PATH)


def build_csv() -> None:
    with CSV_PATH.open("w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f)
        writer.writerow(["단계", "영역", "세부 영역", "항목", "중요도", "담당팀", "기한", "완료 기준", "체크 시점", "리스크", "의존성", "메모"])
        for item in CHECKLIST:
            writer.writerow([item.phase, item.area, item.subarea, item.task, item.priority, item.owner, item.due, item.completion, item.check_point, item.risk, item.dependency, item.notes])


def build_markdown() -> None:
    lines: list[str] = []
    lines.append("# Lucive 오픈베타 런칭 및 운영 계획")
    lines.append("")
    lines.append("## 1. 이번 오픈베타에서 검증할 것")
    lines.append("")
    for idx, v in enumerate(VALIDATION_POINTS, start=1):
        lines.append(f"{idx}. **{v.theme}**: {v.hypothesis}")
        lines.append(f"   KPI: `{v.kpis}`")
        lines.append(f"   성공 기준: `{v.success}`")
        lines.append(f"   오너: `{v.owner}`")
        lines.append("")
    lines.append("## 2. 준비 타임라인")
    lines.append("")
    lines.extend([
        "- `2026-04-23` ~ `2026-04-25`: 오픈베타 목적, 일정, 캠페인 메시지 확정",
        "- `2026-04-26` ~ `2026-04-28`: LP 3개 확정, 세계관 요약, 랜딩 QA, 사전등록 QA, 설문 QA, 온보딩 측정 구조 정리",
        "- `2026-04-29` ~ `2026-04-30`: KR/JP/US 공식 채널 오픈, 커뮤니티 운영 가이드 확정, 1차 스토리라인 5개 완성, 5/1~5/7 콘텐츠 로드맵 준비, 공식사이트 Integration 최종 체크",
        "- `2026-05-01`: 공식 채널 + 유료 캠페인 동시 오픈",
        "- `2026-05-01` ~ `2026-05-07`: 모집 진행, 일자별 개발/커뮤니티 스토리라인 운영, 2차 스토리라인 5개 완성, 퍼널 품질 모니터링",
        "- `2026-05-06`: 모집 마감일을 `5/8`로 갈지 `5/15`로 늘릴지 결정",
        "- `2026-05-08` ~ `2026-05-12`: 선발, 온보딩 가이드, KPI 트래커, 버그 응대 체계, FGI 후속 설문 마무리",
        "- `2026-05-14`: 3차 스토리라인 5개 완성, 전체 리허설과 Go / No-Go 판단",
        "- `2026-05-15` or `2026-05-22`: 권장 오픈베타 런칭일",
        "",
    ])
    lines.append("## 3. 4주 운영 타임라인")
    lines.append("")
    lines.extend([
        "- `Week 1`: 일일 체크업, 버그/피드백 운영 안정화, 1차 FGI 대상 확정",
        "- `Week 2`: D7 초기 지표 체크, 알려진 이슈 보드 운영, 1차/2차 FGI 진행",
        "- `Week 3`: 중간 KPI / 커뮤니티 / 버그 triage 재정렬, 운영 정책 업데이트",
        "- `Week 4`: 종료 커뮤니케이션, 보강 인터뷰, 최종 KPI / FGI 종합, 종료 공지",
        "",
    ])
    lines.append("## 4. 전체 체크리스트")
    lines.append("")
    lines.append("| 단계 | 영역 | 세부 영역 | 항목 | 중요도 | 담당팀 | 기한 | 완료 기준 | 체크 시점 | 리스크 |")
    lines.append("|---|---|---|---|---|---|---|---|---|---|")
    for item in CHECKLIST:
        lines.append(f"| {item.phase} | {item.area} | {item.subarea} | {item.task} | {item.priority} | {item.owner} | {item.due} | {item.completion} | {item.check_point} | {item.risk} |")
    MD_PATH.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    build_pptx()
    build_csv()
    build_markdown()
    print(f"Created: {PPTX_PATH}")
    print(f"Created: {CSV_PATH}")
    print(f"Created: {MD_PATH}")
